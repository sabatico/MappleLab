"""
FP MAppleLab — Executive Presentation Generator
Forcepoint brand-aligned PowerPoint deck.
Run: python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ─────────────────────────────────────────────
# Forcepoint Brand Palette
# ─────────────────────────────────────────────
FP = {
    "teal":         RGBColor(0x00, 0x75, 0x65),  # Forcepoint primary teal #007565
    "dark":         RGBColor(0x1D, 0x25, 0x2C),  # Forcepoint dark #1D252C
    "white":        RGBColor(0xFF, 0xFF, 0xFF),
    "light_gray":   RGBColor(0xF2, 0xF4, 0xF5),  # light bg
    "mid_gray":     RGBColor(0x8A, 0x93, 0x9B),  # secondary text
    "dark_gray":    RGBColor(0x4A, 0x54, 0x5E),  # body text
    "teal_light":   RGBColor(0x00, 0x9E, 0x8A),  # lighter teal accent
    "teal_dark":    RGBColor(0x00, 0x5C, 0x50),  # darker teal
    "accent_blue":  RGBColor(0x00, 0x6D, 0xAA),  # complement blue
    "amber":        RGBColor(0xF5, 0xA6, 0x23),  # warning/highlight
    "red":          RGBColor(0xE8, 0x4D, 0x3D),  # alert
    "green":        RGBColor(0x2E, 0xCC, 0x71),  # success
    "panel_bg":     RGBColor(0x24, 0x2F, 0x39),  # dark card bg
    "border":       RGBColor(0x3A, 0x47, 0x55),  # subtle border
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color):
    sp = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp


def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=Pt(1), radius=0):
    shape_type = 5 if radius else 1  # rounded rect or rect
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = border_width
    else:
        sp.line.fill.background()
    if radius:
        try:
            spPr = sp.element.find(qn("p:spPr"))
            if spPr is None:
                spPr = sp.element
            prstGeom = spPr.find(qn("a:prstGeom"))
            if prstGeom is not None:
                avLst = prstGeom.find(qn("a:avLst"))
                if avLst is not None:
                    for gd in avLst.findall(qn("a:gd")):
                        if gd.get("name") == "adj":
                            gd.set("fmla", f"val {radius}")
        except Exception:
            pass
    return sp


def add_text_box(slide, x, y, w, h, text, font_size=14, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_name="Segoe UI",
                 italic=False, v_anchor=MSO_ANCHOR.TOP):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = v_anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else FP["white"]
    run.font.name = font_name
    return txb


def add_multiline_text(slide, x, y, w, h, lines, font_size=14, line_spacing=1.2,
                       color=None, font_name="Segoe UI", bold=False, align=PP_ALIGN.LEFT):
    """Add textbox with multiple paragraphs."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color if color else FP["white"]
        run.font.name = font_name
    return txb


def add_bullet_list(slide, x, y, w, h, items, font_size=13, color=None,
                    bullet_char="\u25B8", font_name="Segoe UI", spacing=6):
    """Add bulleted list."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = f"  {bullet_char}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color if color else FP["white"]
        run.font.name = font_name
    return txb


def add_icon_bullet_list(slide, x, y, w, items, font_size=13, icon_size=13,
                         color=None, font_name="Segoe UI", row_height=Inches(0.42)):
    """Add list with colored icon markers."""
    for i, (icon, text, icon_color) in enumerate(items):
        iy = y + i * row_height
        # Icon circle
        circle = slide.shapes.add_shape(
            9,  # oval
            x, iy + Inches(0.04), Inches(0.22), Inches(0.22)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = icon_color
        circle.line.fill.background()
        # Icon text
        add_text_box(slide, x + Inches(0.02), iy + Inches(0.02), Inches(0.22), Inches(0.22),
                     icon, icon_size, bold=True, color=FP["dark"],
                     align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        # Label text
        add_text_box(slide, x + Inches(0.32), iy, w - Inches(0.32), row_height,
                     text, font_size, color=color if color else FP["white"])


def fp_title_bar(slide, y=Inches(0)):
    """Top teal accent bar."""
    add_rect(slide, 0, y, SLIDE_W, Inches(0.06), FP["teal"])


def fp_footer(slide, text="FORCEPOINT CONFIDENTIAL"):
    """Footer bar with confidentiality notice."""
    fy = SLIDE_H - Inches(0.36)
    add_rect(slide, 0, fy, SLIDE_W, Inches(0.36), FP["teal_dark"])
    add_text_box(slide, Inches(0.5), fy + Inches(0.04), Inches(12), Inches(0.28),
                 text, 8, color=FP["white"], align=PP_ALIGN.LEFT,
                 font_name="Segoe UI")


def slide_title(slide, title, subtitle=""):
    """Standard title area for content slides."""
    fp_title_bar(slide)
    add_text_box(slide, Inches(0.6), Inches(0.25), Inches(11), Inches(0.6),
                 title, 28, bold=True, color=FP["white"], font_name="Segoe UI")
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.82), Inches(11), Inches(0.4),
                     subtitle, 14, color=FP["teal_light"], font_name="Segoe UI")


def section_divider(slide, number, title, subtitle=""):
    """Section divider slide with large number."""
    fill_bg(slide, FP["teal"])
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(2.5), Inches(3.0),
                 f"{number:02d}", 120, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 font_name="Segoe UI", align=PP_ALIGN.LEFT)
    # Vertical line
    add_rect(slide, Inches(3.8), Inches(2.0), Pt(2), Inches(2.8), FP["white"])
    add_text_box(slide, Inches(4.2), Inches(2.4), Inches(7.5), Inches(1.0),
                 title, 36, bold=True, color=FP["white"], font_name="Segoe UI")
    if subtitle:
        add_text_box(slide, Inches(4.2), Inches(3.4), Inches(7.5), Inches(1.0),
                     subtitle, 16, color=FP["light_gray"], font_name="Segoe UI")
    fp_footer(slide)


def card_box(slide, x, y, w, h, fill=None, border=None):
    """Dark card panel."""
    return add_rect(slide, x, y, w, h,
                    fill if fill else FP["panel_bg"],
                    border_color=border if border else FP["border"],
                    border_width=Pt(1), radius=8000)


# ─────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────

def build_presentation():
    prs = new_prs()

    # ══════════════════════════════════════════
    # SLIDE 1 — TITLE SLIDE
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    # Large teal accent block on left
    add_rect(sl, 0, 0, Inches(0.35), SLIDE_H, FP["teal"])
    fp_title_bar(sl)

    add_text_box(sl, Inches(1.2), Inches(1.5), Inches(10), Inches(0.5),
                 "FP MAppleLab", 48, bold=True, color=FP["teal_light"], font_name="Segoe UI")
    add_text_box(sl, Inches(1.2), Inches(2.2), Inches(10), Inches(1.2),
                 "Leveraging Apple Silicon for On-Demand\nmacOS Access Across Forcepoint",
                 28, color=FP["white"], font_name="Segoe UI")
    add_text_box(sl, Inches(1.2), Inches(3.8), Inches(10), Inches(0.8),
                 "A Corporate Private Cloud for Native macOS Virtual Machines\nRepurposed Hardware  |  BSL 1.1 Licensed  |  One-Time Production License",
                 14, color=FP["mid_gray"], font_name="Segoe UI")

    # Bottom tag line
    add_rect(sl, Inches(1.2), Inches(5.8), Inches(8), Pt(1.5), FP["teal"])
    add_text_box(sl, Inches(1.2), Inches(6.0), Inches(10), Inches(0.4),
                 "Presented to Forcepoint Leadership", 12,
                 color=FP["mid_gray"], font_name="Segoe UI")
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 2 — THE CHALLENGE
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "The Challenge: macOS Access Is Fragmented",
                "Teams need macOS daily — but access is inconsistent and limited")

    # Left: Pain points
    card_box(sl, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.8))
    add_text_box(sl, Inches(0.9), Inches(1.65), Inches(5.3), Inches(0.4),
                 "Today's Situation", 18, bold=True, color=FP["amber"])

    pain_items = [
        ("\u2718", "Only some team members have corporate Macs — access is role-dependent", FP["red"]),
        ("\u2718", "Go4Labs Hackintosh VMs are limited, unstable, and lack native performance", FP["red"]),
        ("\u2718", "Physical Mac pool is scarce — often one device shared per team", FP["red"]),
        ("\u2718", "No self-service provisioning — requests go through manual processes", FP["red"]),
        ("\u2718", "No persistent environments — work lost between sessions", FP["red"]),
    ]
    add_icon_bullet_list(sl, Inches(0.9), Inches(2.2), Inches(5.2),
                         pain_items, font_size=12, color=FP["light_gray"],
                         row_height=Inches(0.52))

    # Right: Who needs it
    card_box(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(4.8))
    add_text_box(sl, Inches(7.1), Inches(1.65), Inches(5.3), Inches(0.4),
                 "Who Needs macOS Access?", 18, bold=True, color=FP["teal_light"])

    teams = [
        ("\u25CF", "Technical Support — reproducing Mac-specific issues", FP["teal"]),
        ("\u25CF", "Professional Services — client environment replication", FP["teal"]),
        ("\u25CF", "Pre-Sales & SE — live demos on native macOS", FP["teal"]),
        ("\u25CF", "QA Engineering — cross-platform test coverage", FP["teal"]),
        ("\u25CF", "Development — build and debug macOS agents", FP["teal"]),
        ("\u25CF", "Lab / POC teams — proof of concept environments", FP["teal"]),
    ]
    add_icon_bullet_list(sl, Inches(7.1), Inches(2.2), Inches(5.2),
                         teams, font_size=12, color=FP["light_gray"],
                         row_height=Inches(0.52))

    # Bottom conclusion
    add_rect(sl, Inches(0.6), Inches(6.5), Inches(12), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.45),
                 "Every team that touches macOS needs reliable, on-demand access — not a shared workaround",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 3 — THE VISION / BEST OUTCOME
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "The Vision: macOS On-Demand for Everyone",
                "A private cloud that puts native Mac VMs at every team's fingertips")

    # Central vision statement
    add_rect(sl, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2), FP["teal"],
             radius=10000)
    add_text_box(sl, Inches(1.8), Inches(1.65), Inches(9.7), Inches(1.1),
                 "Extensive adoption of a corporate private cloud to cover\nall macOS VM needs — for every team, every use case, on demand.",
                 18, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)

    # Three pillars
    pillar_w = Inches(3.2)
    pillar_h = Inches(3.2)
    pillar_y = Inches(3.2)
    pillars = [
        ("Self-Service", "Any authorized user can\nspin up a macOS VM\nin under a second\nfrom a web portal.",
         FP["teal"], "\U0001F680"),
        ("Native Performance", "Apple Silicon M-series\nvirtualization delivers\nfull macOS fidelity —\nnot a Hackintosh.",
         FP["accent_blue"], "\u26A1"),
        ("Minimal Investment", "Leverages decommissioned\nMac Mini & MacBook Pro\nhardware. One-time BSL 1.1\nlicense for production use.",
         FP["green"], "\U0001F4B0"),
    ]
    for i, (title, desc, color, icon) in enumerate(pillars):
        px = Inches(0.8) + i * (pillar_w + Inches(0.3))
        card_box(sl, px, pillar_y, pillar_w, pillar_h, border=color)
        # Colored top bar on card
        add_rect(sl, px, pillar_y, pillar_w, Inches(0.06), color)
        add_text_box(sl, px + Inches(0.1), pillar_y + Inches(0.2), pillar_w - Inches(0.2), Inches(0.45),
                     f"{icon}  {title}", 18, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(sl, px + Inches(0.2), pillar_y + Inches(0.75), pillar_w - Inches(0.4), Inches(2.2),
                     desc, 13, color=FP["light_gray"], align=PP_ALIGN.CENTER)

    # Bottom line
    add_rect(sl, Inches(0.6), Inches(6.5), Inches(12), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.45),
                 "FP MAppleLab turns idle hardware into an always-ready macOS lab for the entire organization",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 4 — INTRODUCING FP MAppleLab
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Introducing FP MAppleLab",
                "How we deliver on-demand macOS VMs across Forcepoint")

    add_text_box(sl, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
                 "FP MAppleLab is a web-based UI and orchestration layer for Apple-native virtualization,\nbuilt on top of open-source TART runtime on Apple Silicon hardware.",
                 14, color=FP["light_gray"])

    # Key components grid  2x2
    comp_w = Inches(5.5)
    comp_h = Inches(1.8)
    comps = [
        ("Software: MAppleLab UI + TART Runtime",
         "MAppleLab is the UI and orchestration layer (Flask web portal)\n"
         "with multi-user auth, VM lifecycle, VNC console, and admin analytics.\n"
         "TART is the open-source VM runtime providing Apple Virtualization.framework integration.",
         FP["teal"]),
        ("Hardware: Decommissioned Macs",
         "Repurposes Forcepoint's existing Mac Mini and MacBook Pro fleet\n"
         "that were already slated for disposal. M1/M2/M3 Apple Silicon\n"
         "nodes join the cluster via a simple agent deployment script.",
         FP["accent_blue"]),
        ("Network: Go4Labs Integration",
         "Once running, each macOS VM has internet access and connects\n"
         "to existing Go4Labs environments via port mapping.\n"
         "Browser or direct VNC access from anywhere on corporate network.",
         FP["amber"]),
        ("Licensing: Apple + BSL 1.1",
         "Apple EULA compliant — native Virtualization.framework,\n"
         "max 2 VMs per physical node. MAppleLab is distributed under\n"
         "BSL 1.1: source-available, production use requires a one-time license.",
         FP["green"]),
    ]
    for i, (title, desc, color) in enumerate(comps):
        cx = Inches(0.6) + (i % 2) * (comp_w + Inches(0.5))
        cy = Inches(2.4) + (i // 2) * (comp_h + Inches(0.25))
        card_box(sl, cx, cy, comp_w, comp_h, border=color)
        add_rect(sl, cx, cy, Inches(0.06), comp_h, color)
        add_text_box(sl, cx + Inches(0.25), cy + Inches(0.12), comp_w - Inches(0.4), Inches(0.4),
                     title, 15, bold=True, color=color)
        add_text_box(sl, cx + Inches(0.25), cy + Inches(0.52), comp_w - Inches(0.4), Inches(1.2),
                     desc, 11, color=FP["light_gray"])

    # Bottom
    add_rect(sl, Inches(0.6), Inches(6.5), Inches(12), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.45),
                 "FP MAppleLab combines proven open-source tools with existing Forcepoint assets",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 5 — KEY BENEFITS
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Key Benefits: Why FP MAppleLab Changes the Game",
                "Tangible advantages for teams and operations")

    benefits = [
        ("\U0001F310", "Go4Labs Integration",
         "Every VM gets internet access and connectivity to existing lab environments via port mapping. Seamless.",
         FP["teal"]),
        ("\u26A1", "Instant Startup",
         "VMs launch in under a second. No waiting, no queuing, no tickets. Instant access to the latest macOS.",
         FP["teal_light"]),
        ("\U0001F4E6", "Golden Images",
         "Pre-configured VMs with Forcepoint policies, agents, and tools pre-installed. Ready to use out of the box.",
         FP["accent_blue"]),
        ("\U0001F4BE", "Long-Term Storage",
         "Save, archive, and resume VMs at any time. Work persists across sessions with Docker Registry backing.",
         FP["green"]),
        ("\U0001F4C8", "Scalability",
         "Add new nodes in minutes with an install script. Scale the cluster as demand grows.",
         FP["amber"]),
        ("\u2696\uFE0F", "Fully Licensed & Compliant",
         "Apple EULA compliant (native Virtualization.framework, 2 VMs/node). MAppleLab under BSL 1.1 — one-time production license.",
         FP["teal"]),
        ("\U0001F3AF", "Native Performance",
         "M-series silicon virtualization is on par with physical Mac. Not emulation. Full macOS fidelity.",
         FP["accent_blue"]),
        ("\U0001F5A5\uFE0F", "Browser & Native VNC",
         "Access from any browser via noVNC, or use native macOS Screen Sharing for full-fidelity experience.",
         FP["teal_light"]),
    ]

    bw = Inches(2.85)
    bh = Inches(1.18)
    for i, (icon, title, desc, color) in enumerate(benefits):
        bx = Inches(0.6) + (i % 4) * (bw + Inches(0.2))
        by = Inches(1.5) + (i // 4) * (bh + Inches(0.2))
        card_box(sl, bx, by, bw, bh, border=color)
        add_rect(sl, bx, by, bw, Inches(0.04), color)
        add_text_box(sl, bx + Inches(0.12), by + Inches(0.1), bw - Inches(0.2), Inches(0.35),
                     f"{icon}  {title}", 12, bold=True, color=color)
        add_text_box(sl, bx + Inches(0.12), by + Inches(0.45), bw - Inches(0.2), Inches(0.65),
                     desc, 10, color=FP["light_gray"])

    add_rect(sl, Inches(0.6), Inches(6.5), Inches(12), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.45),
                 "FP MAppleLab delivers instant, compliant, native macOS access with minimal infrastructure investment",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 6 — MAJOR FUNCTIONALITY
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Platform Capabilities at a Glance",
                "MAppleLab delivers full VM lifecycle management through a unified web portal")

    # User features
    card_box(sl, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.0))
    add_text_box(sl, Inches(0.7), Inches(1.6), Inches(5.5), Inches(0.4),
                 "User Features", 18, bold=True, color=FP["teal_light"])

    user_feats = [
        "Create macOS VMs from curated base images (Sonoma, Sequoia, etc.)",
        "Start / Stop / Delete VMs with one click",
        "Save & Archive VMs to persistent registry storage",
        "Resume archived VMs on any available node",
        "Migrate VMs between nodes seamlessly",
        "Browser-based VNC console (noVNC) — works from any device",
        "Native macOS Screen Sharing (.vncloc download)",
        "Real-time operation progress with stage indicators",
        "Per-user quota management and usage tracking",
    ]
    add_bullet_list(sl, Inches(0.7), Inches(2.1), Inches(5.5), Inches(4.0),
                    user_feats, font_size=11, color=FP["light_gray"],
                    bullet_char="\u25B8", spacing=5)

    # Admin features
    card_box(sl, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.0))
    add_text_box(sl, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.4),
                 "Admin Features", 18, bold=True, color=FP["amber"])

    admin_feats = [
        "Multi-user management with invite-based onboarding",
        "Role-based access control (Admin / User)",
        "Node fleet management — add, activate, deactivate, drain",
        "Real-time node health monitoring (CPU, RAM, Disk, Slots)",
        "Cross-user operations dashboard for all VMs",
        "Usage analytics — VM lifetime & VNC session tracking",
        "Docker Registry storage management & orphan cleanup",
        "SMTP configuration for automated invite emails",
        "Platform settings and operational controls",
    ]
    add_bullet_list(sl, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.0),
                    admin_feats, font_size=11, color=FP["light_gray"],
                    bullet_char="\u25B8", spacing=5)

    add_rect(sl, Inches(0.5), Inches(6.65), Inches(12.2), Inches(0.45), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.7), Inches(6.7), Inches(11.8), Inches(0.35),
                 "Full-featured web portal — from VM creation to admin analytics, all in one place",
                 12, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 7-11 — SCREENSHOT PLACEHOLDER SLIDES
    # ══════════════════════════════════════════
    screenshot_slides = [
        ("Dashboard — My VMs",
         "The main user view showing all VMs with status badges, quick actions, and quota usage.\n"
         "Users see their personal VM inventory with real-time status polling.",
         "INSERT SCREENSHOT: My VMs Dashboard (https://192.168.1.195/)"),
        ("VM Detail & Operations",
         "Detailed VM view with operation progress, transfer stages, and console access buttons.\n"
         "One-click Start, Stop, Save, Resume, Migrate, and Delete actions.",
         "INSERT SCREENSHOT: VM Detail page with running VM"),
        ("VNC Console — Browser Access",
         "Full macOS desktop accessible from any browser via noVNC WebSocket bridge.\n"
         "Supports bandwidth profiles and works over HTTPS from any device.",
         "INSERT SCREENSHOT: noVNC console showing macOS desktop"),
        ("Admin — Nodes & Health",
         "Fleet management view showing all Mac nodes with live CPU, RAM, disk, and VM slot status.\n"
         "Add, activate, deactivate, and drain nodes from this panel.",
         "INSERT SCREENSHOT: Admin Nodes page with health stats"),
        ("Admin — Usage Analytics",
         "Cross-user VM and VNC session analytics with lifetime composition bars.\n"
         "Track how teams use macOS resources across the organization.",
         "INSERT SCREENSHOT: Admin Usage page with analytics"),
    ]

    for title, desc, placeholder in screenshot_slides:
        sl = blank_slide(prs)
        fill_bg(sl, FP["dark"])
        slide_title(sl, title, desc)

        # Screenshot placeholder area
        ph_x, ph_y = Inches(0.8), Inches(1.8)
        ph_w, ph_h = Inches(11.7), Inches(4.8)
        card_box(sl, ph_x, ph_y, ph_w, ph_h,
                 fill=RGBColor(0x15, 0x1C, 0x28), border=FP["teal"])
        add_text_box(sl, ph_x + Inches(1), ph_y + Inches(1.8), ph_w - Inches(2), Inches(1.2),
                     placeholder, 16, color=FP["mid_gray"], align=PP_ALIGN.CENTER,
                     v_anchor=MSO_ANCHOR.MIDDLE)
        # Instruction for presenter
        add_text_box(sl, ph_x + Inches(1), ph_y + Inches(3.2), ph_w - Inches(2), Inches(0.6),
                     "(Right-click this area > Change Picture > select your screenshot)",
                     11, italic=True, color=FP["dark_gray"], align=PP_ALIGN.CENTER)

        fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 12 — DEMO VIDEO PLACEHOLDER
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Live Demo",
                "Watch FP MAppleLab in action — from VM creation to console access")

    # Video placeholder
    ph_x, ph_y = Inches(1.5), Inches(1.8)
    ph_w, ph_h = Inches(10.3), Inches(4.5)
    card_box(sl, ph_x, ph_y, ph_w, ph_h,
             fill=RGBColor(0x10, 0x15, 0x20), border=FP["teal"])
    # Play button triangle
    add_text_box(sl, ph_x + Inches(3.5), ph_y + Inches(1.2), Inches(3.3), Inches(2.0),
                 "\u25B6", 80, bold=True, color=FP["teal"],
                 align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(sl, ph_x + Inches(1), ph_y + Inches(3.5), ph_w - Inches(2), Inches(0.6),
                 "INSERT DEMO VIDEO HERE\n(Insert > Video > select your recorded demo file)",
                 14, italic=True, color=FP["mid_gray"], align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.6), Inches(6.5), Inches(12), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.45),
                 "From zero to macOS desktop in under 10 seconds",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 13 — NETWORK ARCHITECTURE
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Network Architecture",
                "How FP MAppleLab integrates with Forcepoint's corporate network")

    # Network zones
    zones = [
        (Inches(0.4), Inches(3.0), "Corporate Network\n(User Browsers)",
         FP["panel_bg"], FP["teal"], Inches(2.5)),
        (Inches(3.3), Inches(4.5), "Manager Node\n(MAppleLab UI + Orchestrator)",
         FP["panel_bg"], FP["accent_blue"], Inches(3.0)),
        (Inches(6.7), Inches(4.5), "Node Fleet\n(Mac Mini / MBP Cluster)",
         FP["panel_bg"], FP["teal"], Inches(3.8)),
        (Inches(10.9), Inches(2.5), "Shared Services\n(Registry + Go4Labs)",
         FP["panel_bg"], FP["green"], Inches(2.0)),
    ]

    zone_y = Inches(1.5)
    for zx, zh, zt, zfill, zborder, zw in zones:
        card_box(sl, zx, zone_y, zw, zh, fill=zfill, border=zborder)
        add_text_box(sl, zx + Inches(0.1), zone_y + Inches(0.1), zw - Inches(0.2), Inches(0.5),
                     zt.split("\n")[0], 11, bold=True, color=zborder, align=PP_ALIGN.CENTER)
        if "\n" in zt:
            add_text_box(sl, zx + Inches(0.1), zone_y + Inches(0.35), zw - Inches(0.2), Inches(0.3),
                         zt.split("\n")[1], 9, color=FP["mid_gray"], align=PP_ALIGN.CENTER)

    # User zone details
    user_items = ["HTTPS Web Portal", "WSS VNC Console", "vnc:// Screen Sharing"]
    for i, item in enumerate(user_items):
        iy = zone_y + Inches(0.75) + i * Inches(0.5)
        add_rect(sl, Inches(0.6), iy, Inches(2.1), Inches(0.38),
                 FP["panel_bg"], border_color=FP["teal"], radius=5000)
        add_text_box(sl, Inches(0.7), iy + Inches(0.04), Inches(1.9), Inches(0.3),
                     item, 9, color=FP["light_gray"], align=PP_ALIGN.CENTER)

    # Manager zone details
    mgr_items = ["TLS Reverse Proxy\n(Caddy/nginx)", "Flask + Gunicorn\nWeb Application",
                 "WebSocket VNC Bridge", "Direct TCP Proxy\n(57000-57099)"]
    for i, item in enumerate(mgr_items):
        iy = zone_y + Inches(0.75) + i * Inches(0.7)
        add_rect(sl, Inches(3.5), iy, Inches(2.6), Inches(0.58),
                 FP["panel_bg"], border_color=FP["accent_blue"], radius=5000)
        add_text_box(sl, Inches(3.6), iy + Inches(0.04), Inches(2.4), Inches(0.5),
                     item, 9, color=FP["light_gray"], align=PP_ALIGN.CENTER)

    # Node zone details
    node_items = ["TART Agent (:7000)", "TART Runtime (VM)", "websockify VNC (:5900+)",
                  "Local VM Disk Storage"]
    for i, item in enumerate(node_items):
        iy = zone_y + Inches(0.75) + i * Inches(0.7)
        add_rect(sl, Inches(6.9), iy, Inches(3.4), Inches(0.58),
                 FP["panel_bg"], border_color=FP["teal"], radius=5000)
        add_text_box(sl, Inches(7.0), iy + Inches(0.04), Inches(3.2), Inches(0.5),
                     item, 9, color=FP["light_gray"], align=PP_ALIGN.CENTER)

    # Shared services details
    svc_items = ["Docker Registry\n(:5001)", "Go4Labs Env\n(port mapping)"]
    for i, item in enumerate(svc_items):
        iy = zone_y + Inches(0.75) + i * Inches(0.7)
        add_rect(sl, Inches(11.05), iy, Inches(1.7), Inches(0.58),
                 FP["panel_bg"], border_color=FP["green"], radius=5000)
        add_text_box(sl, Inches(11.1), iy + Inches(0.04), Inches(1.55), Inches(0.5),
                     item, 9, color=FP["light_gray"], align=PP_ALIGN.CENTER)

    # Protocol legend
    ly = Inches(5.8)
    legend = [
        (FP["teal"], "HTTPS — user access"),
        (FP["accent_blue"], "HTTPS — management API"),
        (FP["red"], "WSS/TCP — VNC console"),
        (FP["green"], "OCI — VM artefacts"),
    ]
    for i, (lc, lt) in enumerate(legend):
        lx = Inches(0.6) + i * Inches(3.0)
        sp = sl.shapes.add_shape(1, lx, ly, Inches(0.2), Inches(0.12))
        sp.fill.solid()
        sp.fill.fore_color.rgb = lc
        sp.line.fill.background()
        add_text_box(sl, lx + Inches(0.28), ly - Inches(0.02), Inches(2.5), Inches(0.2),
                     lt, 9, color=FP["light_gray"])

    add_rect(sl, Inches(0.5), Inches(6.5), Inches(12.2), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.7), Inches(6.55), Inches(11.8), Inches(0.45),
                 "Seamless corporate network integration with TLS, VNC, and Go4Labs connectivity",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 14 — SOFTWARE ARCHITECTURE
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Software Architecture",
                "Modular Flask application with agent-based node orchestration")

    # Manager core
    mgr_x, mgr_y = Inches(0.5), Inches(1.5)
    mgr_w, mgr_h = Inches(7.5), Inches(3.8)
    card_box(sl, mgr_x, mgr_y, mgr_w, mgr_h, border=FP["accent_blue"])
    add_rect(sl, mgr_x, mgr_y, mgr_w, Inches(0.35), FP["accent_blue"])
    add_text_box(sl, mgr_x + Inches(0.2), mgr_y + Inches(0.04), mgr_w - Inches(0.4), Inches(0.28),
                 "MANAGER — MAppleLab UI & Orchestrator (Flask + Gunicorn)", 11, bold=True,
                 color=FP["dark"], align=PP_ALIGN.CENTER)

    # Module boxes
    mods = [
        ("Auth & Sessions", FP["teal"]),
        ("VM Lifecycle\nOrchestrator", FP["accent_blue"]),
        ("Admin Module\nUser & Node Mgmt", FP["amber"]),
        ("REST / HTMX\nAPI Layer", FP["accent_blue"]),
        ("Node Scheduler\n& Health", FP["teal"]),
        ("Registry Inventory\n& Cleanup", FP["green"]),
        ("VNC WebSocket\nBridge", FP["red"]),
        ("SSH Tunnel\nManager", FP["red"]),
        ("Direct TCP\nProxy (.vncloc)", FP["red"]),
        ("Usage Telemetry\n& Metrics", FP["amber"]),
    ]
    mw, mh = Inches(1.38), Inches(0.62)
    mg = Inches(0.12)
    cols = 5
    for i, (name, col) in enumerate(mods):
        ci = i % cols
        ri = i // cols
        mx = mgr_x + Inches(0.15) + ci * (mw + mg)
        my = mgr_y + Inches(0.5) + ri * (mh + mg)
        add_rect(sl, mx, my, mw, mh, col, border_color=FP["border"], radius=8000)
        add_text_box(sl, mx + Inches(0.05), my + Inches(0.03), mw - Inches(0.1), mh - Inches(0.06),
                     name, 8, bold=True, color=FP["dark"], align=PP_ALIGN.CENTER,
                     v_anchor=MSO_ANCHOR.MIDDLE)

    # Data layer label
    add_text_box(sl, mgr_x + Inches(0.15), mgr_y + mgr_h - Inches(0.4),
                 mgr_w - Inches(0.3), Inches(0.3),
                 "Python 3.10+  |  Flask-SQLAlchemy  |  Paramiko SSH  |  websockify  |  noVNC",
                 9, color=FP["mid_gray"], align=PP_ALIGN.CENTER)

    # Data stores on right
    ds_x = Inches(8.4)
    ds_w = Inches(4.3)
    stores = [
        ("Portal Database (SQLite)", "Users, VMs, Nodes, Settings,\nStatus Events, VNC Sessions", FP["accent_blue"]),
        ("Docker Registry (OCI)", "Archived VM images,\nMigration staging, Quota accounting", FP["green"]),
        ("Application Logs", "File-based operational logging", FP["mid_gray"]),
    ]
    for i, (title, desc, color) in enumerate(stores):
        sy = Inches(1.5) + i * Inches(1.2)
        card_box(sl, ds_x, sy, ds_w, Inches(1.0), border=color)
        add_rect(sl, ds_x, sy, Inches(0.06), Inches(1.0), color)
        add_text_box(sl, ds_x + Inches(0.2), sy + Inches(0.08), ds_w - Inches(0.3), Inches(0.3),
                     title, 11, bold=True, color=color)
        add_text_box(sl, ds_x + Inches(0.2), sy + Inches(0.38), ds_w - Inches(0.3), Inches(0.55),
                     desc, 9, color=FP["light_gray"])

    # Node fleet box
    nf_x, nf_y = Inches(0.5), Inches(5.55)
    nf_w = Inches(7.5)
    nf_h = Inches(1.3)
    card_box(sl, nf_x, nf_y, nf_w, nf_h, border=FP["teal"])
    add_rect(sl, nf_x, nf_y, nf_w, Inches(0.3), FP["teal"])
    add_text_box(sl, nf_x + Inches(0.2), nf_y + Inches(0.02), nf_w - Inches(0.4), Inches(0.25),
                 "CLIENT NODE FLEET (x N)", 10, bold=True, color=FP["dark"], align=PP_ALIGN.CENTER)

    node_mods = [
        ("TART Agent\n(HTTP :7000)", FP["teal"]),
        ("TART Runtime\n(tart CLI)", FP["teal"]),
        ("Local VM\nStorage", FP["panel_bg"]),
        ("websockify\nVNC Bridge", FP["red"]),
        ("VM VNC\n(:5900)", FP["red"]),
    ]
    for i, (nm, nc) in enumerate(node_mods):
        nx = nf_x + Inches(0.2) + i * (Inches(1.38) + Inches(0.12))
        ny = nf_y + Inches(0.4)
        add_rect(sl, nx, ny, Inches(1.38), Inches(0.72), nc, border_color=FP["border"], radius=6000)
        add_text_box(sl, nx + Inches(0.05), ny + Inches(0.04), Inches(1.28), Inches(0.64),
                     nm, 8, bold=True, color=FP["white"] if nc != FP["panel_bg"] else FP["light_gray"],
                     align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    add_rect(sl, Inches(0.5), Inches(6.95), Inches(12.2), Inches(0.2), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.7), Inches(6.95), Inches(11.8), Inches(0.2),
                 "Modular, extensible architecture — each component independently testable and replaceable",
                 10, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 15 — HARDWARE ARCHITECTURE
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Hardware Architecture",
                "Repurposed Apple Silicon devices form a scalable private cloud")

    # Manager node
    card_box(sl, Inches(0.5), Inches(1.5), Inches(3.8), Inches(3.5), border=FP["accent_blue"])
    add_rect(sl, Inches(0.5), Inches(1.5), Inches(3.8), Inches(0.35), FP["accent_blue"])
    add_text_box(sl, Inches(0.6), Inches(1.52), Inches(3.6), Inches(0.3),
                 "MANAGER NODE", 12, bold=True, color=FP["dark"], align=PP_ALIGN.CENTER)

    mgr_specs = [
        "Mac Mini (M1/M2) or MacBook Pro",
        "macOS Sonoma / Sequoia",
        "MAppleLab UI & Orchestrator (Flask)",
        "TLS Reverse Proxy (Caddy)",
        "Docker Registry (:5001)",
        "SQLite Database",
        "SSH key access to all nodes",
    ]
    add_bullet_list(sl, Inches(0.7), Inches(2.0), Inches(3.4), Inches(2.8),
                    mgr_specs, font_size=10, color=FP["light_gray"],
                    bullet_char="\u25B8", spacing=4)

    # Node fleet
    card_box(sl, Inches(4.7), Inches(1.5), Inches(8.1), Inches(3.5), border=FP["teal"])
    add_rect(sl, Inches(4.7), Inches(1.5), Inches(8.1), Inches(0.35), FP["teal"])
    add_text_box(sl, Inches(4.8), Inches(1.52), Inches(7.9), Inches(0.3),
                 "CLIENT NODE FLEET (Expandable)", 12, bold=True, color=FP["dark"], align=PP_ALIGN.CENTER)

    # Individual node cards
    for i in range(3):
        nx = Inches(4.9) + i * Inches(2.6)
        card_box(sl, nx, Inches(2.0), Inches(2.35), Inches(2.8),
                 fill=RGBColor(0x18, 0x22, 0x2E), border=FP["border"])
        node_name = f"Node {chr(65+i)}"
        add_text_box(sl, nx + Inches(0.1), Inches(2.05), Inches(2.15), Inches(0.3),
                     node_name, 11, bold=True, color=FP["teal_light"], align=PP_ALIGN.CENTER)
        node_specs = [
            "Mac Mini M1/M2",
            "8-16 GB RAM",
            "256+ GB SSD",
            "macOS + TART",
            "TART Agent (:7000)",
            "Max 2 VMs (EULA)",
        ]
        add_bullet_list(sl, nx + Inches(0.15), Inches(2.4), Inches(2.05), Inches(2.3),
                        node_specs, font_size=9, color=FP["light_gray"],
                        bullet_char="\u25B8", spacing=3)

    # + More nodes indicator
    add_text_box(sl, Inches(11.7), Inches(2.8), Inches(1.0), Inches(1.0),
                 "+N", 28, bold=True, color=FP["teal"], align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(11.2), Inches(3.6), Inches(1.5), Inches(0.4),
                 "Scale as needed", 9, color=FP["mid_gray"], align=PP_ALIGN.CENTER)

    # Apple compliance badge
    card_box(sl, Inches(0.5), Inches(5.3), Inches(12.2), Inches(0.9), border=FP["green"])
    add_rect(sl, Inches(0.5), Inches(5.3), Inches(0.06), Inches(0.9), FP["green"])
    add_text_box(sl, Inches(0.8), Inches(5.35), Inches(11.6), Inches(0.4),
                 "Apple EULA Compliance: 2 VMs per physical macOS host  |  Apple Virtualization.framework  |  No Hackintosh",
                 12, bold=True, color=FP["green"], align=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(0.8), Inches(5.75), Inches(11.6), Inches(0.35),
                 "Core hardware sourced from Forcepoint's decommissioned Mac inventory — minimal accessories investment (e.g. Thunderbolt-to-RJ45 adapters)",
                 10, color=FP["light_gray"], align=PP_ALIGN.CENTER)

    add_rect(sl, Inches(0.5), Inches(6.5), Inches(12.2), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.7), Inches(6.55), Inches(11.8), Inches(0.45),
                 "Recycled hardware, enterprise-grade capability — FP MAppleLab maximizes asset value",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 16 — BOM / TECHNOLOGY STACK
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Technology Stack & Bill of Materials",
                "Major components powering FP MAppleLab")

    # Category columns
    categories = [
        {
            "title": "Core Platform",
            "color": FP["teal"],
            "items": [
                ("macOS", "Host OS (manager + nodes)"),
                ("Python 3.10+", "Runtime for Flask app"),
                ("Flask", "Web application framework"),
                ("Flask-SQLAlchemy", "ORM / data access"),
                ("SQLite", "Application database"),
            ]
        },
        {
            "title": "VM & Orchestration",
            "color": FP["accent_blue"],
            "items": [
                ("TART", "VM runtime (Cirrus Labs)"),
                ("tart_agent", "Node-side API service"),
                ("Paramiko", "SSH transport / tunnels"),
                ("Flask-Sock", "WebSocket VNC bridge"),
            ]
        },
        {
            "title": "Registry & Storage",
            "color": FP["green"],
            "items": [
                ("Docker CLI", "Container management"),
                ("Colima", "Docker engine on macOS"),
                ("Registry:2", "VM image archive"),
            ]
        },
        {
            "title": "Network & Access",
            "color": FP["amber"],
            "items": [
                ("OpenSSH", "Key auth & remote exec"),
                ("noVNC", "Browser VNC client"),
                ("websockify", "WS-TCP VNC bridge"),
                ("Caddy/nginx", "TLS reverse proxy"),
            ]
        },
        {
            "title": "Security",
            "color": FP["red"],
            "items": [
                ("Flask-Login", "Session auth"),
                ("Flask-Bcrypt", "Password hashing"),
                ("AGENT_TOKEN", "Node API auth"),
                ("SECRET_KEY", "Session signing"),
            ]
        },
    ]

    cat_w = Inches(2.35)
    cat_gap = Inches(0.15)
    cat_y = Inches(1.5)

    for ci, cat in enumerate(categories):
        cx = Inches(0.4) + ci * (cat_w + cat_gap)
        # Category header
        add_rect(sl, cx, cat_y, cat_w, Inches(0.35), cat["color"], radius=4000)
        add_text_box(sl, cx + Inches(0.1), cat_y + Inches(0.04), cat_w - Inches(0.2), Inches(0.28),
                     cat["title"], 11, bold=True, color=FP["dark"], align=PP_ALIGN.CENTER)

        # Items
        for ri, (comp, role) in enumerate(cat["items"]):
            ry = cat_y + Inches(0.45) + ri * Inches(0.72)
            card_box(sl, cx, ry, cat_w, Inches(0.62), border=cat["color"])
            add_text_box(sl, cx + Inches(0.1), ry + Inches(0.04), cat_w - Inches(0.2), Inches(0.25),
                         comp, 10, bold=True, color=cat["color"])
            add_text_box(sl, cx + Inches(0.1), ry + Inches(0.3), cat_w - Inches(0.2), Inches(0.25),
                         role, 9, color=FP["light_gray"])

    add_rect(sl, Inches(0.4), Inches(6.5), Inches(12.4), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.45),
                 "Open-source foundations  |  BSL 1.1 licensed  |  One-time license grants full production use and modification rights",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 17 — ARCHITECTURE DIAGRAM (Env)
    # Embed the Mermaid-based diagrams as reference slides
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Reference: Environment Architecture Diagram",
                "End-to-end user journey — from browser to VM console and archived storage")

    # Placeholder for the generated diagram PPTX
    card_box(sl, Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.0),
             fill=RGBColor(0x12, 0x19, 0x2B), border=FP["teal"])
    add_text_box(sl, Inches(1.5), Inches(3.0), Inches(10.0), Inches(1.5),
                 "EMBED DIAGRAM FROM:\n/docs/schemas/pr/1_environment_architecture_user_experience.pptx\n\n"
                 "(Copy the diagram slide from this file, or regenerate using generate_diagrams.py)",
                 14, color=FP["mid_gray"], align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 18 — ARCHITECTURE DIAGRAM (Technical)
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Reference: Technical Architecture Diagram",
                "Internal modules, data stores, and protocol-level communication channels")

    card_box(sl, Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.0),
             fill=RGBColor(0x12, 0x19, 0x2B), border=FP["accent_blue"])
    add_text_box(sl, Inches(1.5), Inches(3.0), Inches(10.0), Inches(1.5),
                 "EMBED DIAGRAM FROM:\n/docs/schemas/pr/2_technical_architecture_modules_and_channels.pptx\n\n"
                 "(Copy the diagram slide from this file, or regenerate using generate_diagrams.py)",
                 14, color=FP["mid_gray"], align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 19 — ARCHITECTURE DIAGRAM (Portal)
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Reference: Portal Structure Diagram",
                "Component hierarchy, user-facing features, and admin functions")

    card_box(sl, Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.0),
             fill=RGBColor(0x12, 0x19, 0x2B), border=FP["amber"])
    add_text_box(sl, Inches(1.5), Inches(3.0), Inches(10.0), Inches(1.5),
                 "EMBED DIAGRAM FROM:\n/docs/schemas/pr/3_portal_structure_components_and_functions.pptx\n\n"
                 "(Copy the diagram slide from this file, or regenerate using generate_diagrams.py)",
                 14, color=FP["mid_gray"], align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 20 — NEXT STEPS / CALL TO ACTION
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    slide_title(sl, "Next Steps: Scaling FP MAppleLab Across Forcepoint",
                "From pilot to production — a clear path forward")

    # Timeline / steps
    steps = [
        ("01", "Pilot Validation",
         "Current deployment proves concept with real\nusers across TS, PS, and engineering teams.",
         FP["teal"], "COMPLETE"),
        ("02", "Hardware Collection",
         "Identify and collect decommissioned Mac Mini\nand MacBook Pro units from IT inventory.",
         FP["accent_blue"], "IN PROGRESS"),
        ("03", "Cluster Expansion",
         "Deploy TART agent on collected hardware.\nAdd nodes to the cluster via install script.",
         FP["amber"], "NEXT"),
        ("04", "Organization Rollout",
         "Onboard remaining teams. Configure golden\nimages with standard Forcepoint policies.",
         FP["green"], "PLANNED"),
    ]

    sw = Inches(2.8)
    sh = Inches(3.0)
    sy = Inches(1.5)
    for i, (num, title, desc, color, status) in enumerate(steps):
        sx = Inches(0.5) + i * (sw + Inches(0.2))
        card_box(sl, sx, sy, sw, sh, border=color)
        add_rect(sl, sx, sy, sw, Inches(0.06), color)

        # Number circle
        circle = sl.shapes.add_shape(9, sx + Inches(1.1), sy + Inches(0.2), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text_box(sl, sx + Inches(1.1), sy + Inches(0.2), Inches(0.55), Inches(0.55),
                     num, 18, bold=True, color=FP["dark"], align=PP_ALIGN.CENTER,
                     v_anchor=MSO_ANCHOR.MIDDLE)

        add_text_box(sl, sx + Inches(0.15), sy + Inches(0.85), sw - Inches(0.3), Inches(0.35),
                     title, 14, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(sl, sx + Inches(0.15), sy + Inches(1.25), sw - Inches(0.3), Inches(1.0),
                     desc, 10, color=FP["light_gray"], align=PP_ALIGN.CENTER)

        # Status badge
        badge_w = Inches(1.2)
        badge_colors = {
            "COMPLETE": FP["green"], "IN PROGRESS": FP["amber"],
            "NEXT": FP["accent_blue"], "PLANNED": FP["mid_gray"],
        }
        add_rect(sl, sx + (sw - badge_w) / 2, sy + sh - Inches(0.45),
                 badge_w, Inches(0.28), badge_colors[status], radius=14000)
        add_text_box(sl, sx + (sw - badge_w) / 2, sy + sh - Inches(0.45),
                     badge_w, Inches(0.28),
                     status, 8, bold=True, color=FP["dark"], align=PP_ALIGN.CENTER,
                     v_anchor=MSO_ANCHOR.MIDDLE)

    # Ask
    card_box(sl, Inches(0.5), Inches(4.8), Inches(12.2), Inches(1.4), border=FP["teal"])
    add_rect(sl, Inches(0.5), Inches(4.8), Inches(12.2), Inches(0.06), FP["teal"])
    add_text_box(sl, Inches(0.8), Inches(4.95), Inches(11.6), Inches(0.4),
                 "What We Need", 16, bold=True, color=FP["teal_light"])

    asks = [
        "\u25B8  Approval to proceed with hardware collection from IT decommission inventory",
        "\u25B8  Coordination with Go4Labs team for network port mapping integration",
        "\u25B8  Stakeholder buy-in from team leads for onboarding and golden image requirements",
    ]
    add_multiline_text(sl, Inches(0.8), Inches(5.4), Inches(11.6), Inches(0.8),
                       asks, font_size=11, color=FP["light_gray"])

    add_rect(sl, Inches(0.5), Inches(6.5), Inches(12.2), Inches(0.55), FP["teal_dark"],
             border_color=FP["teal"], radius=6000)
    add_text_box(sl, Inches(0.7), Inches(6.55), Inches(11.8), Inches(0.45),
                 "FP MAppleLab is ready to scale — let's bring native macOS access to every team that needs it",
                 13, bold=True, color=FP["white"], align=PP_ALIGN.CENTER,
                 v_anchor=MSO_ANCHOR.MIDDLE)
    fp_footer(sl)

    # ══════════════════════════════════════════
    # SLIDE 21 — CLOSING / THANK YOU
    # ══════════════════════════════════════════
    sl = blank_slide(prs)
    fill_bg(sl, FP["dark"])
    add_rect(sl, 0, 0, Inches(0.35), SLIDE_H, FP["teal"])
    fp_title_bar(sl)

    add_text_box(sl, Inches(1.2), Inches(2.0), Inches(10), Inches(0.8),
                 "FP MAppleLab", 48, bold=True, color=FP["teal_light"])
    add_text_box(sl, Inches(1.2), Inches(2.9), Inches(10), Inches(0.6),
                 "On-Demand macOS. Native Performance. Minimal Investment.", 24,
                 color=FP["white"])

    add_rect(sl, Inches(1.2), Inches(4.0), Inches(8), Pt(1.5), FP["teal"])

    add_text_box(sl, Inches(1.2), Inches(4.4), Inches(10), Inches(0.5),
                 "Thank You", 28, bold=True, color=FP["white"])
    add_text_box(sl, Inches(1.2), Inches(5.0), Inches(10), Inches(0.5),
                 "Questions & Discussion", 16, color=FP["mid_gray"])

    fp_footer(sl)

    return prs


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

if __name__ == "__main__":
    prs = build_presentation()
    out_path = os.path.join(OUT_DIR, "FP_MAppleLab_Executive_Presentation.pptx")
    prs.save(out_path)
    print(f"  Presentation saved to: {out_path}")
    print(f"  Total slides: {len(prs.slides)}")
