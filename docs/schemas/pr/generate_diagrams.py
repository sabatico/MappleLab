"""
MAppleLab — presentation diagram generator.
Produces three standalone .pptx files (one diagram per file) in the same folder.
The generated diagrams include browser VNC, native `.vncloc` direct TCP paths,
and admin usage analytics references.
Run: python generate_diagrams.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ─────────────────────────────────────────────
# Colour palette
# ─────────────────────────────────────────────
C = {
    "bg":            RGBColor(0x12, 0x19, 0x2B),  # deep navy
    "panel":         RGBColor(0x1E, 0x28, 0x40),  # dark panel
    "border":        RGBColor(0x2E, 0x3D, 0x60),  # panel border
    "blue":          RGBColor(0x2D, 0x7D, 0xD2),  # primary blue
    "teal":          RGBColor(0x0D, 0xAF, 0xA1),  # teal / nodes
    "green":         RGBColor(0x06, 0xD6, 0x8A),  # success / archive
    "amber":         RGBColor(0xFF, 0xB7, 0x0E),  # admin / caution
    "purple":        RGBColor(0x8A, 0x56, 0xC9),  # portal / user
    "red":           RGBColor(0xEF, 0x47, 0x6F),  # VNC / console
    "white":         RGBColor(0xFF, 0xFF, 0xFF),
    "light":         RGBColor(0xB0, 0xBE, 0xD4),  # secondary text
    "arrow":         RGBColor(0x4D, 0x6E, 0xA8),  # connector line
    "arrow_bright":  RGBColor(0x5B, 0xCE, 0xFA),  # highlighted connector
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color):
    """Solid background rectangle covering the whole slide."""
    sp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, SLIDE_H
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.zorder = 0
    return sp


def box(slide, x, y, w, h, fill, text="", font_size=11, bold=False,
        text_color=None, border_color=None, border_width=Pt(1),
        radius=0, align=PP_ALIGN.CENTER, v_anchor="middle"):
    """Add a rounded rectangle with optional label."""
    from pptx.util import Pt as _Pt

    # Shape type 5 = ROUNDED_RECTANGLE in MSO auto-shape numbering
    sp = slide.shapes.add_shape(5, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill

    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = border_width
    else:
        sp.line.fill.background()

    # Tune corner radius via XML (gracefully skip if element structure differs)
    try:
        spPr = sp.element.find(qn("p:spPr"))
        if spPr is None:
            spPr = sp.element
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is not None:
            avLst = prstGeom.find(qn("a:avLst"))
            if avLst is not None:
                for gd in avLst.findall(qn("a:gd")):
                    if gd.get("name") == "adj":
                        gd.set("fmla", f"val {radius if radius else 12000}")
    except Exception:
        pass

    if text:
        tf = sp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = _Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color if text_color else C["white"]
        tf.margin_top    = _Pt(4)
        tf.margin_bottom = _Pt(4)
        tf.margin_left   = _Pt(6)
        tf.margin_right  = _Pt(6)

    return sp


def label(slide, x, y, w, h, text, font_size=10, bold=False,
          color=None, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color if color else C["white"]
    return txb


def connector(slide, x1, y1, x2, y2, color=None, width=Pt(1.5),
              label_text="", label_size=8):
    """Draw a straight connector line between two points using the pptx API."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from pptx.oxml.ns import qn as _qn
    from lxml import etree

    clr = color if color else C["arrow"]
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2
    )
    conn.line.color.rgb = clr
    conn.line.width = width

    # Add tail arrowhead via XML (safe — element already exists)
    try:
        ln_el = conn.element.find(".//" + _qn("a:ln"))
        if ln_el is None:
            # Locate spPr and create ln element
            spPr = conn.element.find(_qn("p:spPr"))
            if spPr is None:
                spPr = conn.element
            ln_el = etree.SubElement(spPr, _qn("a:ln"))
        tail = ln_el.find(_qn("a:tailEnd"))
        if tail is None:
            tail = etree.SubElement(ln_el, _qn("a:tailEnd"))
        tail.set("type", "arrow")
        tail.set("w", "med")
        tail.set("len", "med")
    except Exception:
        pass

    if label_text:
        mid_x = (x1 + x2) // 2
        mid_y = (y1 + y2) // 2
        lw = Inches(1.4)
        lh = Inches(0.22)
        txb = slide.shapes.add_textbox(mid_x - lw//2, mid_y - lh//2, lw, lh)
        tf  = txb.text_frame
        p   = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label_text
        run.font.size = Pt(label_size)
        run.font.color.rgb = C["arrow_bright"]
        run.font.italic = True


def section_header(slide, x, y, w, text, color):
    """Coloured section heading bar."""
    sp = slide.shapes.add_shape(1, x, y, w, Inches(0.32))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    tf = sp.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = C["bg"]


def slide_title(slide, text, sub=""):
    label(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.42),
          text, font_size=20, bold=True, color=C["white"])
    if sub:
        label(slide, Inches(0.3), Inches(0.52), Inches(12.7), Inches(0.28),
              sub, font_size=11, color=C["light"])


# ─────────────────────────────────────────────
# Diagram 1 — Environment Architecture
# ─────────────────────────────────────────────

def build_diagram1():
    prs  = new_prs()
    sl   = blank_slide(prs)
    fill_bg(sl, C["bg"])
    slide_title(sl,
                "MAppleLab — Environment Architecture",
                "End-to-end user journey from web browser to virtual machines and archived storage")

    # ── Zones ────────────────────────────────
    # x positions (Inches → Emu done inside box())
    zone_y   = Inches(1.0)
    zone_h   = Inches(5.8)
    pad      = Inches(0.15)

    zone_defs = [
        (Inches(0.18), Inches(2.55), C["purple"],  "Users"),
        (Inches(2.88), Inches(2.05), C["blue"],    "Manager Node"),
        (Inches(5.08), Inches(4.05), C["teal"],    "Node Fleet"),
        (Inches(9.28), Inches(2.35), C["green"],   "Archive Storage"),
    ]
    for zx, zw, zc, zt in zone_defs:
        # translucent panel
        sp = sl.shapes.add_shape(1, zx, zone_y, zw, zone_h)
        sp.fill.solid()
        sp.fill.fore_color.rgb = C["panel"]
        sp.line.color.rgb = zc
        sp.line.width = Pt(1.2)
        section_header(sl, zx, zone_y, zw, zt, zc)

    # ── User Zone ────────────────────────────
    bw, bh = Inches(2.2), Inches(0.62)
    bx = Inches(0.28)

    box(sl, bx, Inches(1.52), bw, bh, C["purple"],
        "👤  End Users\nBrowser Access", 10, border_color=C["border"])
    box(sl, bx, Inches(2.42), bw, bh, C["amber"],
        "🔑  Admin Users\nFull Control", 10, border_color=C["border"])
    box(sl, bx, Inches(4.3), bw, Inches(1.4), C["panel"],
        "Web Portal\n(any browser,\nany device)", 9,
        border_color=C["purple"])

    # ── Manager Node ─────────────────────────
    mx = Inches(2.98)
    box(sl, mx, Inches(1.52), Inches(1.85), bh, C["blue"],
        "🌐  MAppleLab Portal\nFlask + Gunicorn", 9, border_color=C["border"])
    box(sl, mx, Inches(2.42), Inches(1.85), bh, C["blue"],
        "🔒  TLS Reverse Proxy\nCaddy / nginx", 9, border_color=C["border"])
    box(sl, mx, Inches(3.35), Inches(1.85), bh, C["panel"],
        "📊  Portal Database\nSQLite / SQL", 9, border_color=C["blue"])
    box(sl, mx, Inches(4.27), Inches(1.85), bh, C["panel"],
        "📡  Node Scheduler\n& Health Monitor", 9, border_color=C["blue"])
    box(sl, mx, Inches(5.20), Inches(1.85), bh, C["red"],
        "🖥  VNC Bridge\nWebSocket Relay", 9, border_color=C["border"])

    # ── Node Fleet ───────────────────────────
    nx = Inches(5.22)
    nw = Inches(1.15)
    nh = Inches(1.12)
    ng = Inches(0.22)
    ny = [Inches(1.52), Inches(2.84), Inches(4.18)]
    for i, (ny_i, lbl) in enumerate(zip(ny, ["Node A", "Node B", "Node C"])):
        box(sl, nx, ny_i, nw*3+Inches(0.18)*2, nh,
            C["teal"], "", border_color=C["border"])
        label(sl, nx+Inches(0.06), ny_i+Inches(0.04), nw*3+Inches(0.18)*2,
              Inches(0.25), lbl, 9, bold=True, color=C["bg"])
        box(sl, nx+Inches(0.06),            ny_i+Inches(0.3),  nw, Inches(0.68),
            C["panel"], "▶ Running\n   VMs", 8, border_color=C["teal"])
        box(sl, nx+nw+Inches(0.22),         ny_i+Inches(0.3),  nw, Inches(0.68),
            C["panel"], "⏸ Stopped\n   VMs", 8, border_color=C["teal"])
        box(sl, nx+nw*2+Inches(0.38),       ny_i+Inches(0.3),  nw, Inches(0.68),
            C["red"],   "🖥 VNC\n   Server", 8, border_color=C["border"])

    box(sl, nx, Inches(5.42), nw*3+Inches(0.18)*2, Inches(0.7),
        C["panel"], "Local TART VM Storage (per node disk)", 9,
        border_color=C["teal"])

    # ── Archive Zone ─────────────────────────
    ax = Inches(9.38)
    aw = Inches(2.1)
    box(sl, ax, Inches(1.52), aw, Inches(1.0), C["green"],
        "🗄  Docker Registry\nVM Image Archive", 10, bold=True,
        border_color=C["border"])
    box(sl, ax, Inches(2.78), aw, Inches(0.75), C["panel"],
        "Archived VMs\n(save & resume)", 9, border_color=C["green"])
    box(sl, ax, Inches(3.73), aw, Inches(0.75), C["panel"],
        "Migration Staging\n(push → pull)", 9, border_color=C["green"])
    box(sl, ax, Inches(4.68), aw, Inches(0.75), C["panel"],
        "Per-user Quota\nDisk Accounting", 9, border_color=C["green"])
    box(sl, ax, Inches(5.63), aw, Inches(0.62), C["panel"],
        "OCI-compatible\ntart push / pull", 9, border_color=C["green"])

    # ── Connectors ───────────────────────────
    # Users → Manager
    connector(sl, Inches(2.48), Inches(1.83), Inches(2.98), Inches(1.83),
              C["purple"], Pt(1.8), "HTTPS")
    # Admin → Manager
    connector(sl, Inches(2.48), Inches(2.73), Inches(2.98), Inches(2.73),
              C["amber"], Pt(1.8), "HTTPS")
    # Manager → Nodes
    connector(sl, Inches(4.83), Inches(1.83), Inches(5.22), Inches(1.83),
              C["blue"], Pt(1.8), "HTTPS API")
    connector(sl, Inches(4.83), Inches(2.73), Inches(5.22), Inches(3.08),
              C["blue"], Pt(1.5))
    connector(sl, Inches(4.83), Inches(3.66), Inches(5.22), Inches(4.42),
              C["blue"], Pt(1.5))
    # Nodes → Archive
    connector(sl, Inches(8.82), Inches(2.08), Inches(9.38), Inches(2.02),
              C["green"], Pt(1.8), "tart push/pull")
    connector(sl, Inches(8.82), Inches(3.40), Inches(9.38), Inches(3.16),
              C["green"], Pt(1.5))
    connector(sl, Inches(8.82), Inches(4.74), Inches(9.38), Inches(4.30),
              C["green"], Pt(1.5))
    # VNC
    connector(sl, Inches(2.48), Inches(4.60), Inches(2.98), Inches(5.51),
              C["red"], Pt(1.8), "WSS VNC")
    connector(sl, Inches(4.83), Inches(5.51), Inches(8.18), Inches(3.50),
              C["red"], Pt(1.5), "WS/SSH Tunnel")

    # ── Legend ───────────────────────────────
    legend_items = [
        (C["purple"], "User access (HTTPS)"),
        (C["blue"],   "Management API (HTTPS)"),
        (C["green"],  "VM artefacts (tart push/pull)"),
        (C["red"],    "VNC console (WSS / SSH Tunnel)"),
    ]
    lx, ly = Inches(0.28), Inches(6.75)
    label(sl, lx, ly - Inches(0.22), Inches(4), Inches(0.2),
          "Legend:", 8, bold=True, color=C["light"])
    for i, (lc, lt) in enumerate(legend_items):
        ix = lx + Inches(i * 3.1)
        sp = sl.shapes.add_shape(1, ix, ly, Inches(0.25), Inches(0.14))
        sp.fill.solid(); sp.fill.fore_color.rgb = lc; sp.line.fill.background()
        label(sl, ix + Inches(0.3), ly - Inches(0.02),
              Inches(2.7), Inches(0.2), lt, 8, color=C["light"])

    return prs


# ─────────────────────────────────────────────
# Diagram 2 — Technical Architecture
# ─────────────────────────────────────────────

def build_diagram2():
    prs = new_prs()
    sl  = blank_slide(prs)
    fill_bg(sl, C["bg"])
    slide_title(sl,
                "MAppleLab — Technical Architecture",
                "Internal modules, data stores, and protocol-level communication channels")

    # ── Manager block ────────────────────────
    mgr_x, mgr_y = Inches(2.55), Inches(1.02)
    mgr_w, mgr_h = Inches(5.6), Inches(5.72)
    sp = sl.shapes.add_shape(1, mgr_x, mgr_y, mgr_w, mgr_h)
    sp.fill.solid(); sp.fill.fore_color.rgb = C["panel"]
    sp.line.color.rgb = C["blue"]; sp.line.width = Pt(1.5)
    section_header(sl, mgr_x, mgr_y, mgr_w, "MANAGER NODE", C["blue"])

    # Module grid inside manager
    mw, mh = Inches(1.65), Inches(0.6)
    mg = Inches(0.12)
    modules = [
        ("Auth & Sessions",    C["purple"]),
        ("VM Lifecycle\nOrchestrator", C["blue"]),
        ("Admin Module\nUser & Node Mgmt", C["amber"]),
        ("API Layer\nHTMX / REST",  C["blue"]),
        ("Node Scheduler\n& Health",   C["teal"]),
        ("Registry\nInventory & Cleanup", C["green"]),
        ("Console\nWebSocket Bridge",  C["red"]),
        ("SSH Tunnel\nManager",    C["red"]),
    ]
    for i, (name, col) in enumerate(modules):
        col_i = i % 3
        row_i = i // 3
        mx = mgr_x + Inches(0.18) + col_i * (mw + mg)
        my = mgr_y + Inches(0.48) + row_i * (mh + mg)
        box(sl, mx, my, mw, mh, col, name, 8,
            border_color=C["border"], radius=12000)

    # Flask label
    label(sl, mgr_x + Inches(0.1), mgr_y + mgr_h - Inches(0.38),
          mgr_w - Inches(0.2), Inches(0.28),
          "Flask / Gunicorn  ·  Python 3.10+", 8,
          color=C["light"])

    # ── Data stores ──────────────────────────
    ds_x = Inches(8.42)
    box(sl, ds_x, Inches(1.42), Inches(2.1), Inches(0.9),
        C["panel"], "📊 Portal Database\nSQLite / PostgreSQL\n(Users · VMs · Nodes · Settings)",
        8, border_color=C["blue"])
    box(sl, ds_x, Inches(2.55), Inches(2.1), Inches(1.05),
        C["panel"], "🗄 Docker Registry\nOCI-compatible\nVM Image Archive\n(tart push / pull)",
        8, border_color=C["green"])
    box(sl, ds_x, Inches(3.82), Inches(2.1), Inches(0.72),
        C["panel"], "📝 Application Logs\nFile-based", 8, border_color=C["border"])

    # ── Ingress block ────────────────────────
    ig_x = Inches(0.18)
    box(sl, ig_x, Inches(1.42), Inches(2.1), Inches(0.72),
        C["purple"], "👤 User Browsers\n(any device)", 9, border_color=C["border"])
    box(sl, ig_x, Inches(2.34), Inches(2.1), Inches(0.72),
        C["amber"],  "🔑 Admin Browsers\n(full access)", 9, border_color=C["border"])
    box(sl, ig_x, Inches(3.26), Inches(2.1), Inches(0.78),
        C["blue"],   "🔒 TLS Reverse Proxy\nCaddy / nginx\n(terminates TLS)", 8,
        border_color=C["border"])

    # ── Node block ───────────────────────────
    nf_x, nf_y = Inches(2.55), Inches(5.0)
    nf_w = Inches(5.6)
    sp2 = sl.shapes.add_shape(1, nf_x, nf_y, nf_w, Inches(2.26))
    sp2.fill.solid(); sp2.fill.fore_color.rgb = C["panel"]
    sp2.line.color.rgb = C["teal"]; sp2.line.width = Pt(1.5)
    section_header(sl, nf_x, nf_y, nf_w, "CLIENT NODE FLEET (×N)", C["teal"])

    node_items = [
        ("TART Agent\n(HTTP :7000)", C["teal"]),
        ("TART Runtime\n(tart CLI)", C["teal"]),
        ("Local VM Storage\n(node disk)", C["panel"]),
        ("websockify\nVNC Bridge", C["red"]),
        ("VM VNC\n(:5900)", C["red"]),
    ]
    nw, nh2 = Inches(0.95), Inches(0.72)
    for i, (nm, nc) in enumerate(node_items):
        nx2 = nf_x + Inches(0.22) + i * (nw + Inches(0.17))
        box(sl, nx2, nf_y + Inches(0.44), nw, nh2, nc, nm, 8,
            border_color=C["border"], radius=10000)

    # ── Connectors with protocol labels ──────
    # Browsers → Proxy
    connector(sl, Inches(2.28), Inches(1.78), Inches(2.55), Inches(1.78),
              C["purple"], Pt(1.8), "HTTPS")
    connector(sl, Inches(2.28), Inches(2.70), Inches(2.55), Inches(2.70),
              C["amber"],  Pt(1.8), "HTTPS")
    # Proxy → Portal
    connector(sl, Inches(2.28), Inches(3.58), Inches(2.55), Inches(3.58),
              C["blue"], Pt(1.8), "HTTP upstream")

    # WS from browser
    connector(sl, Inches(1.23), Inches(1.78), Inches(1.23), Inches(6.88),
              C["red"], Pt(1.5))
    label(sl, Inches(0.18), Inches(4.18), Inches(0.92), Inches(0.6),
          "WSS\n/console/ws\n(VNC)", 7, color=C["red"])

    # Manager → DB
    connector(sl, Inches(8.15), Inches(1.88), Inches(8.42), Inches(1.88),
              C["blue"], Pt(1.5), "SQL ORM")
    # Manager → Registry
    connector(sl, Inches(8.15), Inches(2.90), Inches(8.42), Inches(2.90),
              C["green"], Pt(1.5), "Docker API")
    # Manager → Nodes
    connector(sl, Inches(5.75), Inches(6.74), Inches(5.75), Inches(7.26),
              C["blue"], Pt(1.8), "HTTPS Agent Token")
    # WS Bridge → Node VNC
    connector(sl, Inches(6.28), Inches(6.74), Inches(7.78), Inches(7.13),
              C["red"], Pt(1.5), "WS / SSH Tunnel")
    # Registry ↔ Nodes
    connector(sl, Inches(8.42), Inches(3.08), Inches(8.18), Inches(7.13),
              C["green"], Pt(1.5), "tart push/pull")

    # ── Protocol legend ───────────────────────
    lg = [
        (C["purple"],       "HTTPS — browser access"),
        (C["blue"],         "HTTPS — management API"),
        (C["red"],          "WSS / SSH tunnel — VNC console"),
        (C["green"],        "OCI registry — VM artefacts"),
        (C["arrow_bright"], "SQL ORM — portal data"),
    ]
    lx, ly = Inches(0.18), Inches(7.03)
    for i, (lc, lt) in enumerate(lg):
        ix = lx + Inches(i * 2.52)
        sp = sl.shapes.add_shape(1, ix, ly, Inches(0.22), Inches(0.12))
        sp.fill.solid(); sp.fill.fore_color.rgb = lc; sp.line.fill.background()
        label(sl, ix + Inches(0.28), ly - Inches(0.02),
              Inches(2.2), Inches(0.18), lt, 7, color=C["light"])

    return prs


# ─────────────────────────────────────────────
# Diagram 3 — Portal Structure
# ─────────────────────────────────────────────

def build_diagram3():
    prs = new_prs()
    sl  = blank_slide(prs)
    fill_bg(sl, C["bg"])
    slide_title(sl,
                "MAppleLab — Portal Structure",
                "Component hierarchy, user-facing features, and admin functions")

    # ── Root portal box ──────────────────────
    root_w, root_h = Inches(2.5), Inches(0.56)
    root_x = (SLIDE_W - root_w) / 2
    root_y = Inches(1.0)
    box(sl, root_x, root_y, root_w, root_h, C["blue"],
        "🌐  MAppleLab Portal", 12, bold=True, border_color=C["border"])

    # ── Auth layer ───────────────────────────
    auth_y = Inches(1.85)
    auth_w = Inches(2.8)
    auth_x = (SLIDE_W - auth_w) / 2
    box(sl, auth_x, auth_y, auth_w, Inches(0.52), C["purple"],
        "🔒  Auth  ·  Login  ·  Invite Password Setup  ·  Session Guard",
        9, border_color=C["border"])
    connector(sl, SLIDE_W//2, root_y + root_h, SLIDE_W//2, auth_y,
              C["blue"], Pt(1.8))

    # ── Split line ───────────────────────────
    split_y = Inches(2.6)
    connector(sl, Inches(0.5), split_y, Inches(12.8), split_y,
              C["border"], Pt(0.8))

    # Column headings
    col_heads = [
        (Inches(0.25), Inches(4.6), C["purple"], "User Frontend"),
        (Inches(7.0),  Inches(5.7), C["amber"],  "Admin Backend"),
    ]
    for cx, cw, cc, ct in col_heads:
        box(sl, cx, split_y + Inches(0.06), cw, Inches(0.38),
            cc, ct, 11, bold=True, border_color=C["border"])

    # ── User Frontend columns ─────────────────
    uf_y = split_y + Inches(0.62)
    uf_cols = [
        {
            "title": "📋  My VMs Dashboard",
            "color": C["purple"],
            "items": [
                ("VM Status Badges", C["panel"]),
                ("Auto-Refresh Polling", C["panel"]),
                ("Quota Usage Display", C["panel"]),
                ("VM Row Actions", C["panel"]),
            ]
        },
        {
            "title": "⚙️  VM Operations",
            "color": C["blue"],
            "items": [
                ("Create VM", C["panel"]),
                ("Start / Stop", C["panel"]),
                ("Save & Archive", C["panel"]),
                ("Resume from Archive", C["panel"]),
                ("Migrate Between Nodes", C["panel"]),
                ("Delete VM", C["panel"]),
            ]
        },
        {
            "title": "🖥  VNC Console",
            "color": C["red"],
            "items": [
                ("In-Browser noVNC Client", C["panel"]),
                ("Bandwidth/Render Profiles", C["panel"]),
                ("Same-Origin WSS Bridge", C["panel"]),
                ("SSH Tunnel Mode (WAN)", C["panel"]),
            ]
        },
        {
            "title": "📊  VM Detail & Progress",
            "color": C["teal"],
            "items": [
                ("Live Operation Stage", C["panel"]),
                ("Transfer Progress %", C["panel"]),
                ("Raw TART Output Line", C["panel"]),
                ("Cleanup Status Badge", C["panel"]),
            ]
        },
    ]

    uf_col_w = Inches(1.58)
    uf_gap   = Inches(0.12)
    uf_start = Inches(0.25)
    item_h   = Inches(0.33)

    for ci, col in enumerate(uf_cols):
        cx2 = uf_start + ci * (uf_col_w + uf_gap)
        # Column header
        box(sl, cx2, uf_y, uf_col_w, Inches(0.44),
            col["color"], col["title"], 8, bold=True, border_color=C["border"])
        # Items
        for ri, (item_text, item_col) in enumerate(col["items"]):
            iy = uf_y + Inches(0.54) + ri * (item_h + Inches(0.06))
            box(sl, cx2, iy, uf_col_w, item_h, item_col,
                item_text, 8, border_color=col["color"])

    # ── Admin Backend columns ──────────────────
    ab_cols = [
        {
            "title": "👤  User Management",
            "color": C["amber"],
            "items": [
                ("Create Users (Invite)", C["panel"]),
                ("Set Role: Admin/User", C["panel"]),
                ("Quota Controls", C["panel"]),
                ("Resend Invites", C["panel"]),
                ("Delete Users", C["panel"]),
            ]
        },
        {
            "title": "🖧  Node Management",
            "color": C["teal"],
            "items": [
                ("Add/Remove Nodes", C["panel"]),
                ("Node Health View", C["panel"]),
                ("Activate / Deactivate", C["panel"]),
                ("Deactivation Drain", C["panel"]),
            ]
        },
        {
            "title": "🗄  Registry Storage",
            "color": C["green"],
            "items": [
                ("Storage Usage Bar", C["panel"]),
                ("Trackable Artefacts", C["panel"]),
                ("Orphan Detection", C["panel"]),
                ("Delete Orphans", C["panel"]),
                ("Retry Cleanup", C["panel"]),
            ]
        },
        {
            "title": "🔧  Platform Settings",
            "color": C["blue"],
            "items": [
                ("SMTP Configuration", C["panel"]),
                ("Send Test Email", C["panel"]),
                ("Ops Dashboard", C["panel"]),
                ("Cross-user VM Actions", C["panel"]),
            ]
        },
    ]

    ab_start = Inches(7.0)
    for ci, col in enumerate(ab_cols):
        cx2 = ab_start + ci * (uf_col_w + uf_gap)
        box(sl, cx2, uf_y, uf_col_w, Inches(0.44),
            col["color"], col["title"], 8, bold=True, border_color=C["border"])
        for ri, (item_text, item_col) in enumerate(col["items"]):
            iy = uf_y + Inches(0.54) + ri * (item_h + Inches(0.06))
            box(sl, cx2, iy, uf_col_w, item_h, item_col,
                item_text, 8, border_color=col["color"])

    # Auth → columns
    connector(sl, Inches(3.5),  auth_y + Inches(0.52), Inches(3.5),  split_y + Inches(0.06),
              C["purple"], Pt(1.5))
    connector(sl, Inches(9.8),  auth_y + Inches(0.52), Inches(9.8),  split_y + Inches(0.06),
              C["amber"],  Pt(1.5))

    return prs


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

if __name__ == "__main__":
    diagrams = [
        ("1_environment_architecture_user_experience.pptx",  build_diagram1),
        ("2_technical_architecture_modules_and_channels.pptx", build_diagram2),
        ("3_portal_structure_components_and_functions.pptx",  build_diagram3),
    ]
    for fname, builder in diagrams:
        prs = builder()
        out = os.path.join(OUT_DIR, fname)
        prs.save(out)
        print(f"  ✓  {fname}")

    print("\nAll diagrams saved to", OUT_DIR)
